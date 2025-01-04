from pandas import read_excel
# import matplotlib as mpl
# import matplotlib.pyplot as plt
from numpy import timedelta64 as np_timedelta64

import logging

# import argparse
from math import floor

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt

from pathlib import Path
from datetime import datetime

logger = logging.getLogger(__name__)

# def create_parser():
#     parser = argparse.ArgumentParser()
#     parser.add_argument( '-i', '--input', help='path to the input file', type=str, required=True)
#     parser.add_argument( '-o', '--output', help='path to the output directory', type=str, required=True)
#     parser.add_argument( '-n', '--num_weeks', help='number of weeks of data to show, default 1', type=int, default=1)
#     parser.add_argument( '-x', '--xp_top_n', help='show top n students in the xp boost list, default 10', type=int, default=10)
#     parser.add_argument( '-l', '--il_top_n', help='show top n students in the independent learning list, default 10', type=int, default=10)
#     parser.add_argument( '-m', '--il_min_time_mins', help='minimum time in minutes for independent learning list, default 20 mins', type=int, default=20)
#     return parser

# constants
HEADER_ROW = 3

BY_STUDENT = 'By student'
BY_CLASS = 'By class'
BY_YEAR_GROUP = 'By year group'
BY_REGISTRATION_GROUP = 'By registration group'

FIRST_NAME = 'First Name'
SURNAME = 'Surname'
MATHS_CLASS = 'Maths class'
CLASS = 'Class'
YEAR_GROUP = 'Year group'
REGISTRATION_GROUP = 'Reg. Group'

FONT_NAME = 'Calibri'
FONT_SIZE = Pt(24)

Y7 = "Year 7"
Y8 = "Year 8"
Y9 = "Year 9"
Y10 = "Year 10"
Y11 = "Year 11"

C_OT = 'C (OT)'
C_AT = 'C (AT)'
XPB = 'XPB'
TARGET='T'
XP='XP'
IL='IL (h:mm)'


#mpl.style.available

YG_SORTER = {Y7: 0, Y8: 1, Y9: 2, Y10:3, Y11:4}
SORTED_YEAR_GROUPS = [Y7, Y8, Y9, Y10, Y11]

# plt.style.use('bmh')
# #set globally 
# plt.rcParams['axes.labelsize'] = 24


class Parameters(object):
    input_file = ''
    output_dir = ''
    num_weeks = 1
    xp_top_n = 10
    il_top_n = 10
    il_min_time_mins = 20
    process_year_group = True
    process_maths_class = True
    process_reg_group = True
    process_xp_boost = True
    process_independent_learning = True

    def __str__(self):
        return str(self.__class__) + ": " + str(self.__dict__)


class DataToPublish(object):
    # figures_by_year_group = ''
    # figures_by_maths_class = {}
    # figures_by_reg_group = {}
    yg_data = {}
    rg_data = {}
    mc_data = {}
    xp_data = {}
    independent_learning_data = {}

    def __str__(self):
        return str(self.__class__) + ": " + str(self.__dict__)


# def create_figure(data, xlabel, title, out_dir):
#     fig = data.plot(figsize=(30,15), x=xlabel, kind='bar', title=title, xlabel=xlabel).get_figure()

#     plt.gca().yaxis.set_major_formatter(mpl.ticker.PercentFormatter(xmax=1.0))
#     plt.gca().set_ylim([0, 1])
#     plt.gca().set_yticks(np.arange(0, 1.0, 0.05))
#     plt.gca().grid(True, axis='y')

#     plt.gca().axes.title.set_size(32)
#     plt.gca().axes.tick_params(axis='both', labelsize=20 )

#     fig_path = f'{out_dir}/{title}.png'
#     fig.savefig(fig_path)
#     plt.close(fig)
#     return fig_path


def load_data(input_file, sheet_name):
    # read sheet data in
    data = read_excel(input_file, header=HEADER_ROW, sheet_name=sheet_name, na_values=['-'])
    return data.fillna(0)


def process_individual_students(data_to_publish, params):
    df = load_data(params.input_file, BY_STUDENT)
    if params.process_xp_boost:
        xp_list = get_top_xp_boost(df, params.xp_top_n)
        data_to_publish.xp_data = xp_list.to_dict('list')

    if params.process_independent_learning:
        il_list = get_top_independent_learning(df, params.il_top_n, params.il_min_time_mins)
        data_to_publish.independent_learning_data = il_list.to_dict('list')

    # print(xp_list.to_string(index=False))
    # print(il_list.to_string(index=False))


def get_top_xp_boost(df, top_n):
    df_xp = df[[FIRST_NAME, SURNAME, XP]]
    return df_xp.sort_values(XP, ascending=False).head(top_n)


def get_top_independent_learning(df, top_n, il_time_mins):
    df_il = df[[FIRST_NAME, SURNAME, IL]]
    df_il = df_il[ df_il[IL] > np_timedelta64(il_time_mins, 'm') ]
    return df_il.sort_values(IL, ascending=False).head(top_n)


def process_by_registration_group(data_to_publish, input_file, out_dir, num_weeks):
    df = load_data(input_file, BY_REGISTRATION_GROUP)
    year_groups = set(df[YEAR_GROUP].tolist())
    for yg in year_groups:
        df_yg = df.loc[df[YEAR_GROUP] == yg]
        cols = get_column_names_to_plot(df_yg.columns.tolist(), REGISTRATION_GROUP, num_weeks)
    
        # fig_path = create_figure(df_yg[cols], xlabel=REGISTRATION_GROUP, title=f'{yg} Completion by Form Group', out_dir=out_dir)
        # data_to_publish.figures_by_reg_group[yg] = fig_path

        data_to_publish.rg_data[yg] = df_yg[cols].to_dict('list')


def process_by_maths_class(data_to_publish, input_file, out_dir, num_weeks):
    df = load_data(input_file, BY_CLASS)
    year_groups = set(df[YEAR_GROUP].tolist())
    for yg in year_groups:
        df_yg = df.loc[df[YEAR_GROUP] == yg]
        cols = get_column_names_to_plot(df_yg.columns.tolist(), MATHS_CLASS, num_weeks)
    
        # fig_path = create_figure(df_yg[cols], xlabel=MATHS_CLASS, title=f'{yg} Completion by Maths Class', out_dir=out_dir)
        # data_to_publish.figures_by_maths_class[yg] = fig_path
        data_to_publish.mc_data[yg] = df_yg[cols].to_dict('list')


def get_column_names_to_plot(df_indices, xlabel, num_weeks):
    col_names = [xlabel, C_OT]
    for i in range(1, num_weeks):
        col_name = f'{C_OT}.{i}'
        if col_name in df_indices:
            col_names.append(col_name)
        else:
            logger.debug(f'{col_name} not found in data frame column names. num_weeks: {num_weeks}, dataframe columns: {df_indices}')

    return col_names


def process_by_year_group(data_to_publish, input_file, out_dir, num_weeks):
    df = load_data(input_file, BY_YEAR_GROUP)
    df_sorted = df.sort_values(by=[YEAR_GROUP], key=lambda x: x.map(YG_SORTER))
    cols = get_column_names_to_plot(df_sorted.columns.tolist(), YEAR_GROUP, num_weeks)

    # fig_path = create_figure(df_sorted[cols], xlabel=YEAR_GROUP, title='Completion by Year Group', out_dir=out_dir)
    # data_to_publish.figures_by_year_group = fig_path
    data_to_publish.yg_data = df_sorted[cols].to_dict('list')


def create_output_dir(out_dir):
    today = datetime.today().strftime('%Y%m%d')
    output_path = Path(out_dir, f'sparx_{today}').resolve()
    output_path.mkdir(parents=True, exist_ok=True)
    return output_path


def create_slides(output_path, data_to_publish, params):
    prs = Presentation()
    prs.save(f'{output_path}/sparx_leaderboard.pptx')
    title_slide_layout = prs.slide_layouts[0]
    blank_slide = prs.slide_layouts[5]

    # title page
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title.text_frame.paragraphs[0]
    add_title(title, f'Sparx Leaderboard - {datetime.today().strftime('%Y/%m/%d')}')

    if params.process_year_group:
        # completion by year group
        slide = prs.slides.add_slide(blank_slide)
        title = slide.shapes.title.text_frame.paragraphs[0]
        add_title(title, f'Which year group was the best?')
        add_chart(slide, data_to_publish.yg_data, title)

    # add_image(slide, data_to_publish.figures_by_year_group)

    if params.process_reg_group:
        # completion by registration group
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title.text_frame.paragraphs[0]
        add_title(title, f'Completion by Form Class')
        for year_group in SORTED_YEAR_GROUPS:
            slide = prs.slides.add_slide(blank_slide)
            title = slide.shapes.title.text_frame.paragraphs[0]
            add_title(title, year_group)
            add_chart(slide, data_to_publish.rg_data[year_group], title)


    if params.process_maths_class:
        # completion by maths class
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title.text_frame.paragraphs[0]
        add_title(title, f'Completion by Maths Class')
        for year_group in SORTED_YEAR_GROUPS:
            slide = prs.slides.add_slide(blank_slide)
            title = slide.shapes.title.text_frame.paragraphs[0]
            # image_path = data_to_publish.figures_by_maths_class[year_group]
            add_title(title, year_group)
            add_chart(slide, data_to_publish.mc_data[year_group], title)
            # add_image(slide, image_path)

    if params.process_xp_boost:
        # xp boost top 10
        slide = prs.slides.add_slide(blank_slide)
        title_para = slide.shapes.title.text_frame.paragraphs[0]
        add_title(title_para, f'XP Boost Top {params.xp_top_n} Leaderboard')
        add_table(slide, data_to_publish.xp_data)

    if params.process_independent_learning:
        # independent learning top 10
        slide = prs.slides.add_slide(blank_slide)
        title_para = slide.shapes.title.text_frame.paragraphs[0]
        add_title(title_para, f'Independent Learning Leaderboard: Top {params.il_top_n} students with over {params.il_min_time_mins} minutes independent learning')
        add_table(slide, data_to_publish.independent_learning_data)

    prs.save(f'{output_path}/sparx_leaderboard.pptx')


def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell


def add_title(title_para, title_text):
    title_para.text = title_text
    title_para.font.name = FONT_NAME
    title_para.font.size = FONT_SIZE


def add_image(slide, image_path):
    pic = slide.shapes.add_picture(image_path, Inches(0.2), Inches(1.2), width=Inches(12), height=Inches(6))


def get_num_of_rows_and_cols(d):
    _, v = next(iter(d.items())) # fetch first item in the table_data dict

    rows = len(v)
    cols = len(d.keys())
    return rows, cols


def to_percentages(values):
    return [ v*100 for v in values ]


def add_chart(slide, data, title):
    # define chart data
    chart_data = CategoryChartData()
    for index, key in enumerate(data):
        if index == 0:
            chart_data.categories = data.get(key)
        else:
            chart_data.add_series(key, to_percentages(data.get(key)))

    # add chart to slide
    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(9), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    value_axis = chart.value_axis
    value_axis.maximum_scale = 100.0

    tick_labels = value_axis.tick_labels
    tick_labels.number_format = '0"%"'
    
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False


def add_table(slide, table_data):
    rows, cols = get_num_of_rows_and_cols(table_data)
    max_width = 28
    width = floor(28 / cols)
    diff = max_width - (width * cols)

    x, y, cx, cy = Inches(diff/2), Inches(1.5), Inches(width), Inches(1.5)
    shape = slide.shapes.add_table(rows=rows+1, cols=cols, left=x, top=y, width=cx, height=cy)
    table = shape.table

    for i, col_name in enumerate(table_data):
        table.cell(0, i).text = col_name

    for row in range(0, rows):
        for col, col_name in enumerate(table_data):
            text = str(table_data.get(col_name)[row])
            table.cell(row+1, col).text = text


def run(args):
    input_file = args.input_file
    num_weeks = args.num_weeks

    out_dir = create_output_dir(args.output_dir)
    data_to_publish = DataToPublish()

    if args.process_year_group:
        process_by_year_group(data_to_publish, input_file, out_dir, num_weeks)
    
    if args.process_maths_class:
        process_by_maths_class(data_to_publish, input_file, out_dir, num_weeks)
    
    if args.process_reg_group:
        process_by_registration_group(data_to_publish, input_file, out_dir, num_weeks)
    
    if args.process_xp_boost or args.process_independent_learning:
        process_individual_students(data_to_publish, args)

    create_slides(out_dir, data_to_publish, args)



if __name__ == "__main__":
    args = Parameters()
    args.num_weeks = 5
    args.input_file = '/Users/THM_1/Documents/jupyter/data.xlsx'
    args.output_dir = '/Users/THM_1/Documents/python_projects/sparx'
    run(args)









